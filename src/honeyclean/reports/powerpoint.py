"""
Generate professional PowerPoint presentations from profiling results.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from datetime import datetime
from pathlib import Path
from typing import Dict, Any
import io
import pandas as pd
import logging

# Set up logger
logger = logging.getLogger(__name__)

from ..config import HoneyCleanConfig
from ..visualizations.generators import VisualizationGenerator
from ..utils.formatters import StatisticalFormatter

class PowerPointGenerator:
    """Generate professional PowerPoint presentations from profiling results."""
    
    def __init__(self, config: HoneyCleanConfig):
        self.config = config
        self.viz_generator = VisualizationGenerator(config)
    
    def create_presentation(self, profiling_results: Dict[str, Any], 
                          df_original) -> str:
        """Create comprehensive PowerPoint presentation with embedded plots."""
        prs = Presentation()
        
        # Set slide dimensions
        prs.slide_width = Inches(self.config.slide_width)
        prs.slide_height = Inches(self.config.slide_height)
        
        # Create slides
        logger.info("📄 Generating title slide...")
        self._create_title_slide(prs, profiling_results)
        logger.info("📊 Generating dataset overview slide...")
        self._create_overview_slide(prs, profiling_results)
        
        # Add data conversion results slide if available
        if 'data_conversion' in profiling_results:
            logger.info("🔄 Generating data type conversion analysis slide...")
            self._create_data_conversion_slide(prs, profiling_results)
        
        # Add enhanced analysis slides
        if 'target_correlation' in profiling_results:
            logger.info("📈 Generating target correlation analysis slide...")
            self._create_target_correlation_slide(prs, profiling_results)
        if 'target_distribution' in profiling_results:
            logger.info("📊 Generating target distribution analysis slide...")
            self._create_target_distribution_slide(prs, profiling_results)
        if 'id_uniqueness' in profiling_results:
            logger.info("🔍 Generating ID uniqueness validation slide...")
            self._create_id_uniqueness_slide(prs, profiling_results)
        
        logger.info("📊 Generating missing values analysis slide...")
        self._create_missing_values_slide(prs, df_original)
        
        # Create two slides per column: visualization + enhanced statistics
        for column_name, column_analysis in profiling_results['columns'].items():
            logger.info(f"📊 Generating column analysis slide for {column_name}...")
            # Slide 1: Enhanced visualization with correlation plots and heatmaps
            self._create_column_slide(prs, column_name, column_analysis, df_original, profiling_results)
            # Slide 2: Enhanced statistics with beautiful formatting
            self._create_column_statistics_slide(prs, column_name, column_analysis)
        
        self._create_enhanced_recommendations_slide(prs, profiling_results)
        
        # Save presentation
        output_path = f"{self.config.output_reports}/data_profiling_report.pptx"
        prs.save(output_path)
        
        return output_path
    
    def _create_title_slide(self, prs: Presentation, results: Dict[str, Any]):
        """Create title slide."""
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = "数据分析报告"
        subtitle.text = f"数据集: {results['dataset_info']['name']}\n生成时间: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
        
        title.text_frame.paragraphs[0].font.size = Pt(44)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    def _create_overview_slide(self, prs: Presentation, results: Dict[str, Any]):
        """Create dataset overview slide."""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "数据集概览"
        
        info = results['dataset_info']
        overview_text = f"""数据集摘要:
- 数据规模: {info['shape'][0]:,} 行 × {info['shape'][1]} 列
- 内存使用: {info['memory_usage_mb']:.1f} MB
- 缺失值: {info['total_missing']:,} ({info['missing_percentage']:.1f}%)
- 重复行: {info['duplicate_count']:,}

列类型分布:
- 数值型: {info['numeric_columns']} 列
- 分类型: {info['categorical_columns']} 列
- 日期型: {info['datetime_columns']} 列
- 文本型: {info['text_columns']} 列"""
        
        content.text = overview_text
    
    def _create_missing_values_slide(self, prs: Presentation, df):
        """Create missing values analysis slide."""
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "缺失值分析"
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        
        # Generate and add plot
        try:
            plot_buffer = self.viz_generator.create_missing_values_summary(df)
            if plot_buffer:
                slide.shapes.add_picture(plot_buffer, Inches(1), Inches(1.5), Inches(11), Inches(5))
        except Exception as e:
            # Add error message if plot fails
            error_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(2))
            error_frame = error_box.text_frame
            error_frame.text = f"图表生成失败: {str(e)}"
    
    def _create_column_slide(self, prs: Presentation, column_name: str, 
                           analysis: Dict[str, Any], df_original, profiling_results: Dict[str, Any]):
        """Create individual column analysis slide with enhanced visualizations."""
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = f"特征分析: {column_name} (Enhanced Feature Analysis)"
        title_frame.paragraphs[0].font.size = Pt(24)
        title_frame.paragraphs[0].font.bold = True
        
        series = df_original[column_name]
        
        # Handle different analysis types with enhanced visualizations
        if analysis.get('type') == 'numeric' and 'error' not in analysis:
            try:
                # Create distribution plot (top left) and percentiles plot (top right)
                hist_buffer, percentile_buffer = self.viz_generator.create_enhanced_numeric_plot_for_ppt(series, column_name)
                
                if hist_buffer:
                    slide.shapes.add_picture(hist_buffer, Inches(0.5), Inches(1.2), Inches(5.8), Inches(2.8))
                
                if percentile_buffer:
                    slide.shapes.add_picture(percentile_buffer, Inches(6.7), Inches(1.2), Inches(5.8), Inches(2.8))
                
                # Create correlation plot with target (bottom left) if target exists
                target_correlation_text = ""
                if 'target_correlation' in profiling_results and self.config.target_col:
                    target_cols = self.config.target_col if isinstance(self.config.target_col, list) else [self.config.target_col]
                    
                    for target_col in target_cols:
                        if column_name != target_col and target_col in df_original.columns:
                            try:
                                corr_buffer = self.viz_generator.create_correlation_plot_for_ppt(
                                    df_original, column_name, target_col)
                                if corr_buffer:
                                    slide.shapes.add_picture(corr_buffer, Inches(0.5), Inches(4.2), Inches(5.8), Inches(2.5))
                                
                                # Get correlation value for display
                                if target_col in profiling_results['target_correlation'] and column_name in profiling_results['target_correlation'][target_col]['correlations']:
                                    corr_val = profiling_results['target_correlation'][target_col]['correlations'][column_name]
                                    target_correlation_text = f"与目标变量 {target_col} 的相关性: {corr_val:.4f}"
                                break
                            except Exception as e:
                                # Add correlation error message
                                corr_error_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(5.8), Inches(2.5))
                                corr_error_frame = corr_error_box.text_frame
                                corr_error_frame.text = f"相关性图表生成失败: {str(e)}"
                                break
                
                # Create heatmap showing top correlated features (bottom right)
                try:
                    heatmap_buffer = self.viz_generator.create_feature_correlation_heatmap_for_ppt(
                        df_original, column_name, top_n=5)
                    if heatmap_buffer:
                        slide.shapes.add_picture(heatmap_buffer, Inches(6.7), Inches(4.2), Inches(5.8), Inches(2.5))
                except Exception as e:
                    # Add heatmap error message
                    heatmap_error_box = slide.shapes.add_textbox(Inches(6.7), Inches(4.2), Inches(5.8), Inches(2.5))
                    heatmap_error_frame = heatmap_error_box.text_frame
                    heatmap_error_frame.text = f"特征相关性热图生成失败: {str(e)}"
                
                # Add correlation info text at very bottom
                if target_correlation_text:
                    corr_info_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.9), Inches(12), Inches(0.4))
                    corr_info_frame = corr_info_box.text_frame
                    corr_info_frame.text = target_correlation_text
                    corr_info_frame.paragraphs[0].font.size = Pt(12)
                    corr_info_frame.paragraphs[0].font.bold = True
                
            except Exception as e:
                # Add error message if plots fail
                error_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(3))
                error_frame = error_box.text_frame
                error_frame.text = f"数值型变量图表生成失败: {str(e)}"
                
        elif analysis.get('type') == 'categorical' and 'error' not in analysis:
            try:
                # Create categorical distribution plot (top)
                plot_buffer = self.viz_generator.create_enhanced_categorical_plot_for_ppt(series, column_name)
                if plot_buffer:
                    slide.shapes.add_picture(plot_buffer, Inches(0.5), Inches(1.2), Inches(12), Inches(3))
                
                # Create categorical correlation analysis (bottom) if target exists
                if self.config.target_col:
                    target_cols = self.config.target_col if isinstance(self.config.target_col, list) else [self.config.target_col]
                    
                    for target_col in target_cols:
                        if target_col in df_original.columns:
                            try:
                                cat_analysis_buffer = self.viz_generator.create_categorical_target_analysis_for_ppt(
                                    df_original, column_name, target_col)
                                if cat_analysis_buffer:
                                    slide.shapes.add_picture(cat_analysis_buffer, Inches(0.5), Inches(4.5), Inches(12), Inches(2.2))
                                break
                            except Exception as e:
                                # Add categorical analysis error message
                                cat_error_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(12), Inches(2.2))
                                cat_error_frame = cat_error_box.text_frame
                                cat_error_frame.text = f"分类变量目标分析失败: {str(e)}"
                                break
                
            except Exception as e:
                # Add error message if plots fail
                error_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(3))
                error_frame = error_box.text_frame
                error_frame.text = f"分类型变量图表生成失败: {str(e)}"
        else:
            # Error case or unsupported type
            error_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(4))
            error_frame = error_box.text_frame
            error_text = f"""变量类型: {analysis.get('type', 'unknown')}
总数: {analysis.get('count', 0):,}
缺失值: {analysis.get('missing_count', 0):,} ({analysis.get('missing_percentage', 0):.1f}%)

{analysis.get('error', '该变量类型暂不支持详细分析')}"""
            error_frame.text = error_text
    
    def _create_recommendations_slide(self, prs: Presentation, results: Dict[str, Any]):
        """Create general recommendations slide."""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        content = slide.placeholders[1]
        
        title.text = "数据清洗建议"
        
        all_recommendations = results.get('general_recommendations', [])
        
        if all_recommendations:
            recommendations_text = "整体数据集建议:\n\n" + "\n".join([f"• {rec}" for rec in all_recommendations])
            content.text = recommendations_text
        else:
            content.text = "数据集状况良好，无特殊建议。"
    
    def _create_target_correlation_slide(self, prs: Presentation, results: Dict[str, Any]):
        """Create target correlation analysis slide."""
        slide_layout = prs.slide_layouts[6]  # 使用空白布局
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "目标变量相关性分析 (Target Correlation Analysis)"
        title_frame.paragraphs[0].font.size = Pt(24)
        title_frame.paragraphs[0].font.bold = True
        
        for target_col, correlations in results['target_correlation'].items():
            # 过滤掉nan值并排序
            valid_correlations = {k: v for k, v in correlations['correlations'].items() 
                                if not pd.isna(v) and k != target_col}
            
            sorted_corrs = sorted(valid_correlations.items(), 
                                key=lambda x: abs(x[1]), reverse=True)
            
            # 获取前20和后20
            top_20 = sorted_corrs[:20]
            bottom_20 = sorted_corrs[-20:] if len(sorted_corrs) > 20 else []
            
            # 左侧 - 前20个最强相关性
            left_content = f"🎯 目标变量: {target_col}\n\n前20个最强相关特征:\n"
            for col, corr_val in top_20:
                strength = StatisticalFormatter._interpret_correlation(abs(corr_val))
                left_content += f"• {col}: {corr_val:.4f} ({strength})\n"
            
            left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(5.5), Inches(6))
            left_frame = left_box.text_frame
            left_frame.word_wrap = True
            left_frame.text = left_content
            
            # 右侧 - 后20个最弱相关性
            if bottom_20:
                right_content = f"\n后20个最弱相关特征:\n"
                for col, corr_val in bottom_20:
                    strength = StatisticalFormatter._interpret_correlation(abs(corr_val))
                    right_content += f"• {col}: {corr_val:.4f} ({strength})\n"
            else:
                right_content = "\n总特征数少于40个，\n无需显示最弱相关特征"
            
            right_box = slide.shapes.add_textbox(Inches(6.5), Inches(1.2), Inches(5.5), Inches(6))
            right_frame = right_box.text_frame
            right_frame.word_wrap = True
            right_frame.text = right_content
            
            # 设置字体样式
            for frame in [left_frame, right_frame]:
                for paragraph in frame.paragraphs:
                    paragraph.font.size = Pt(11)
                    if "🎯" in paragraph.text or "前20个" in paragraph.text or "后20个" in paragraph.text:
                        paragraph.font.bold = True
                        paragraph.font.size = Pt(13)
    
    def _create_target_distribution_slide(self, prs: Presentation, results: Dict[str, Any]):
        """Create target distribution analysis slide."""
        slide_layout = prs.slide_layouts[6]  # 使用空白布局
        slide = prs.slides.add_slide(slide_layout)
        
        # 添加标题
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "目标变量分布分析 (Target Distribution Analysis)"
        title_frame.paragraphs[0].font.size = Pt(24)
        title_frame.paragraphs[0].font.bold = True
        
        # Create content from target distribution
        distribution_content = ""
        for target_col, target_stats in results['target_distribution'].items():
            distribution_content += f"\n🎯 目标变量: {target_col}\n"
            distribution_content += f"类型: {target_stats.get('type', 'Unknown')}\n\n"
            
            if target_stats.get('type') == 'categorical':
                distribution_content += "类别分布 (Category Distribution):\n"
                value_counts = target_stats.get('value_counts', {})
                value_percentages = target_stats.get('value_percentages', {})
                
                for value, count in list(value_counts.items())[:8]:  # Top 8
                    percentage = value_percentages.get(value, 0)
                    distribution_content += f"• {value}: {count:,} ({percentage:.1f}%)\n"
                
                # Class balance info
                class_balance = target_stats.get('class_balance', {})
                if class_balance:
                    largest_pct = class_balance.get('largest_class_percentage', 0)
                    is_balanced = class_balance.get('is_balanced', True)
                    balance_status = "平衡" if is_balanced else "不平衡"
                    distribution_content += f"\n类别平衡: {balance_status} (最大类别: {largest_pct:.1f}%)\n"
            
            elif target_stats.get('type') == 'numeric':
                distribution_content += "数值统计 (Numeric Statistics):\n"
                distribution_content += f"• 平均值: {target_stats.get('mean', 0):.2f}\n"
                distribution_content += f"• 中位数: {target_stats.get('median', 0):.2f}\n"
                distribution_content += f"• 标准差: {target_stats.get('std', 0):.2f}\n"
                distribution_content += f"• 最小值: {target_stats.get('min', 0):.2f}\n"
                distribution_content += f"• 最大值: {target_stats.get('max', 0):.2f}\n"
        
        # Add text box
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(6))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        content_frame.text = distribution_content
        
        # Style the text
        for paragraph in content_frame.paragraphs:
            paragraph.font.size = Pt(12)
            if "🎯" in paragraph.text:
                paragraph.font.bold = True
                paragraph.font.size = Pt(14)
            elif any(keyword in paragraph.text for keyword in ["类别分布", "数值统计", "类别平衡"]):
                paragraph.font.bold = True
                paragraph.font.size = Pt(13)
    
    def _create_id_uniqueness_slide(self, prs: Presentation, results: Dict[str, Any]):
        """Create ID uniqueness validation slide."""
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        
        title = slide.shapes.title
        title.text = "ID列唯一性验证 (ID Uniqueness Validation)"
        
        # Create content from ID uniqueness results
        uniqueness_content = "🆔 ID列唯一性检查结果:\n\n"
        
        id_results = results['id_uniqueness']
        for col, stats in id_results.items():
            if 'error' in stats:
                uniqueness_content += f"• {col}: 错误 - {stats['error']}\n"
                continue
                
            total_count = stats.get('total_count', 0)
            unique_count = stats.get('unique_count', 0)
            duplicate_count = stats.get('duplicate_count', 0)
            uniqueness_pct = stats.get('uniqueness_percentage', 0)
            
            status = "✅ 唯一" if duplicate_count == 0 else f"❌ {duplicate_count:,} 重复"
            
            uniqueness_content += f"• {col}:\n"
            uniqueness_content += f"  - 总数: {total_count:,}\n"
            uniqueness_content += f"  - 唯一: {unique_count:,}\n"
            uniqueness_content += f"  - 唯一性: {uniqueness_pct:.1f}%\n"
            uniqueness_content += f"  - 状态: {status}\n\n"
            
            # Show duplicate examples if any
            duplicate_values = stats.get('duplicate_values', [])
            if duplicate_values:
                uniqueness_content += f"  - 重复值示例: {', '.join(map(str, duplicate_values[:5]))}\n\n"
        
        # Check composite ID if exists
        if 'composite_id_uniqueness' in results:
            composite_stats = results['composite_id_uniqueness']
            if 'error' not in composite_stats:
                composite_cols = composite_stats.get('composite_columns', [])
                composite_unique_pct = composite_stats.get('uniqueness_percentage', 0)
                composite_duplicates = composite_stats.get('duplicate_count', 0)
                
                uniqueness_content += f"🔗 复合ID ({'+'.join(composite_cols)}):\n"
                uniqueness_content += f"  - 唯一性: {composite_unique_pct:.1f}%\n"
                uniqueness_content += f"  - 重复: {composite_duplicates:,}\n"
        
        # Add text box
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(5.5))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        content_frame.text = uniqueness_content
        
        # Style the text
        for paragraph in content_frame.paragraphs:
            paragraph.font.size = Pt(11)
            if "🆔" in paragraph.text or "🔗" in paragraph.text:
                paragraph.font.bold = True
                paragraph.font.size = Pt(14)
    
    def _create_enhanced_recommendations_slide(self, prs: Presentation, results: Dict[str, Any]):
        """Create enhanced recommendations slide with comprehensive suggestions."""
        slide_layout = prs.slide_layouts[6]  # 使用空白布局
        slide = prs.slides.add_slide(slide_layout)

        # 添加标题框（手动）
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "综合数据清洗建议 (Comprehensive Cleaning Recommendations)"
        title_frame.paragraphs[0].font.size = Pt(20)
        title_frame.paragraphs[0].font.bold = True

        recommendations_content = ""

        # General recommendations
        general_recs = results.get('general_recommendations', [])
        if general_recs:
            recommendations_content += "🔧 整体建议 (General Recommendations):\n"
            for i, rec in enumerate(general_recs, 1):
                recommendations_content += f"{i}. {rec}\n"
            recommendations_content += "\n"

        # Target-specific recommendations
        if 'target_correlation' in results:
            recommendations_content += "🎯 目标变量建议 (Target Analysis Recommendations):\n"
            recommendations_content += "• 重点关注高相关性特征进行特征工程\n"
            recommendations_content += "• 检查相关性异常值和潜在数据泄露\n"
            recommendations_content += "• 考虑特征选择以避免多重共线性\n\n"

        # ID-specific recommendations
        if 'id_uniqueness' in results:
            id_issues = []
            for col, stats in results['id_uniqueness'].items():
                if stats.get('duplicate_count', 0) > 0:
                    id_issues.append(col)

            if id_issues:
                recommendations_content += "🆔 ID列建议 (ID Column Recommendations):\n"
                for col in id_issues:
                    recommendations_content += f"• {col}: 存在重复值，需要数据去重或验证\n"
                recommendations_content += "\n"

        # Column-specific priority recommendations
        high_priority_issues = []
        for col_name, analysis in results.get('columns', {}).items():
            missing_pct = analysis.get('missing_percentage', 0)
            if missing_pct > 50:
                high_priority_issues.append(f"{col_name}: 缺失值过高 ({missing_pct:.1f}%)")
            elif analysis.get('zscore_outliers', 0) > analysis.get('count', 0) * 0.1:
                high_priority_issues.append(f"{col_name}: 异常值较多")

        if high_priority_issues:
            recommendations_content += "⚠️ 优先处理问题 (Priority Issues):\n"
            for issue in high_priority_issues[:20]:
                recommendations_content += f"• {issue}\n"

        # 添加内容框
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12), Inches(6))
        content_frame = content_box.text_frame
        content_frame.word_wrap = True
        content_frame.text = recommendations_content if recommendations_content else "数据集状况良好，无特殊清洗建议。"

        # 设置样式
        for paragraph in content_frame.paragraphs:
            paragraph.font.size = Pt(11)
            if any(emoji in paragraph.text for emoji in ["🔧", "🎯", "🆔", "⚠️"]):
                paragraph.font.bold = True
                paragraph.font.size = Pt(13)

    
    def _convert_table_to_text(self, formatted_stats: str) -> str:
        """Convert tabulated statistics to presentation-friendly text."""
        lines = formatted_stats.split('\n')
        presentation_text = ""
        
        for line in lines:
            # Remove table formatting characters but keep content
            if line.strip() and not line.strip().startswith('+') and not line.strip().startswith('|==='):
                # Convert table rows to bullet points
                if '|' in line:
                    parts = [part.strip() for part in line.split('|') if part.strip()]
                    if len(parts) >= 2 and not parts[0].startswith('-'):
                        presentation_text += f"• {parts[0]}: {parts[1]}\n"
                else:
                    # Keep headers and section titles
                    presentation_text += line + "\n"
        
        return presentation_text
    
    def _create_column_statistics_slide(self, prs: Presentation, column_name: str, analysis: Dict[str, Any]):
        """Create beautiful formatted statistics slide for a column."""
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title with icon
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        title_frame = title_box.text_frame
        data_type = analysis.get('type', 'unknown')
        type_icon = {"numeric": "📊", "categorical": "🏷️", "datetime": "📅"}.get(data_type, "📋")
        title_frame.text = f"{type_icon} 统计详情: {column_name} (Statistics Details)"
        title_frame.paragraphs[0].font.size = Pt(24)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
        
        if data_type == 'numeric':
            self._create_numeric_stats_layout(slide, analysis)
        elif data_type == 'categorical':
            self._create_categorical_stats_layout(slide, analysis)
        elif data_type == 'datetime':
            self._create_datetime_stats_layout(slide, analysis)
        else:
            self._create_unknown_stats_layout(slide, analysis)
    
    def _create_numeric_stats_layout(self, slide, analysis: Dict[str, Any]):
        """Create beautiful layout for numeric statistics."""
        # Basic Stats Box (Left Top)
        basic_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(5.8), Inches(2.2))
        basic_frame = basic_box.text_frame
        basic_frame.word_wrap = True
        
        count_val = analysis.get('count', 0)
        missing_count_val = analysis.get('missing_count', 0)
        count_str = f"{count_val:,}" if isinstance(count_val, (int, float)) else str(count_val)
        missing_str = f"{missing_count_val:,}" if isinstance(missing_count_val, (int, float)) else str(missing_count_val)
        
        basic_stats = f"""📈 基础统计 (Basic Statistics)

计数 (Count): {count_str}
缺失 (Missing): {missing_str} ({analysis.get('missing_percentage', 0):.1f}%)
平均数 (Mean): {analysis.get('mean', 0):.4f}
中位数 (Median): {analysis.get('median', 0):.4f}
众数 (Mode): {analysis.get('mode', 'N/A')}
标准差 (Std Dev): {analysis.get('std', 0):.4f}"""
        
        basic_frame.text = basic_stats
        basic_frame.paragraphs[0].font.size = Pt(14)
        basic_frame.paragraphs[0].font.bold = True
        for i in range(1, len(basic_frame.paragraphs)):
            basic_frame.paragraphs[i].font.size = Pt(11)
        
        # Distribution Stats Box (Right Top)
        dist_box = slide.shapes.add_textbox(Inches(6.7), Inches(1.3), Inches(5.8), Inches(2.2))
        dist_frame = dist_box.text_frame
        dist_frame.word_wrap = True
        
        dist_stats = f"""📊 分布统计 (Distribution)

最小值 (Min): {analysis.get('min', 0):.4f}
最大值 (Max): {analysis.get('max', 0):.4f}
范围 (Range): {analysis.get('range', 0):.4f}
Q1 (25%): {analysis.get('q1', 0):.4f}
Q3 (75%): {analysis.get('q3', 0):.4f}
IQR: {analysis.get('iqr', 0):.4f}
偏度 (Skewness): {analysis.get('skewness', 0):.4f}
峰度 (Kurtosis): {analysis.get('kurtosis', 0):.4f}"""
        
        dist_frame.text = dist_stats
        dist_frame.paragraphs[0].font.size = Pt(14)
        dist_frame.paragraphs[0].font.bold = True
        for i in range(1, len(dist_frame.paragraphs)):
            dist_frame.paragraphs[i].font.size = Pt(11)
        
        # Outlier Detection Box (Left Bottom)
        outlier_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(5.8), Inches(1.8))
        outlier_frame = outlier_box.text_frame
        outlier_frame.word_wrap = True
        
        zscore_outliers = analysis.get('zscore_outliers', 0)
        iqr_outliers = analysis.get('iqr_outliers', 0)
        modified_zscore_outliers = analysis.get('modified_zscore_outliers', 0)
        
        zscore_str = f"{zscore_outliers:,}" if isinstance(zscore_outliers, (int, float)) else str(zscore_outliers)
        iqr_str = f"{iqr_outliers:,}" if isinstance(iqr_outliers, (int, float)) else str(iqr_outliers)
        modified_str = f"{modified_zscore_outliers:,}" if isinstance(modified_zscore_outliers, (int, float)) else str(modified_zscore_outliers)
        
        outlier_stats = f"""🔍 异常值检测 (Outlier Detection)

Z-Score异常值: {zscore_str} ({analysis.get('zscore_percentage', 0):.1f}%)
IQR异常值: {iqr_str} ({analysis.get('iqr_percentage', 0):.1f}%)
修正Z-Score: {modified_str} ({analysis.get('modified_zscore_percentage', 0):.1f}%)"""
        
        outlier_frame.text = outlier_stats
        outlier_frame.paragraphs[0].font.size = Pt(14)
        outlier_frame.paragraphs[0].font.bold = True
        for i in range(1, len(outlier_frame.paragraphs)):
            outlier_frame.paragraphs[i].font.size = Pt(11)
        
        # Distribution Tests Box (Right Bottom)
        if analysis.get('shapiro_p_value') is not None:
            test_box = slide.shapes.add_textbox(Inches(6.7), Inches(3.8), Inches(5.8), Inches(1.8))
            test_frame = test_box.text_frame
            test_frame.word_wrap = True
            
            test_stats = f"""🧪 分布检验 (Distribution Tests)

正态性检验 (Normality):
Shapiro-Wilk p值: {analysis.get('shapiro_p_value', 0):.6f}
是否正态分布: {'是 (Yes)' if analysis.get('is_normal', False) else '否 (No)'}"""
            
            if analysis.get('dagostino_p_value') is not None:
                test_stats += f"\nD'Agostino p值: {analysis.get('dagostino_p_value', 0):.6f}"
            
            test_frame.text = test_stats
            test_frame.paragraphs[0].font.size = Pt(14)
            test_frame.paragraphs[0].font.bold = True
            for i in range(1, len(test_frame.paragraphs)):
                test_frame.paragraphs[i].font.size = Pt(11)
        
        # Recommendations Box (Bottom)
        recommendations = analysis.get('recommendations', [])
        if recommendations:
            rec_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.0), Inches(12), Inches(1.3))
            rec_frame = rec_box.text_frame
            rec_frame.word_wrap = True
            
            rec_text = "💡 清洗建议 (Recommendations):\n"
            for i, rec in enumerate(recommendations[:4], 1):  # Show top 4
                rec_text += f"{i}. {rec}\n"
            
            rec_frame.text = rec_text
            rec_frame.paragraphs[0].font.size = Pt(14)
            rec_frame.paragraphs[0].font.bold = True
            for i in range(1, len(rec_frame.paragraphs)):
                rec_frame.paragraphs[i].font.size = Pt(10)
    
    def _create_categorical_stats_layout(self, slide, analysis: Dict[str, Any]):
        """Create beautiful layout for categorical statistics."""
        # Basic Stats Box (Left)
        basic_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(5.8), Inches(2.5))
        basic_frame = basic_box.text_frame
        basic_frame.word_wrap = True
        
        count_val = analysis.get('count', 0)
        missing_count_val = analysis.get('missing_count', 0)
        unique_count_val = analysis.get('unique_count', 0)
        
        count_str = f"{count_val:,}" if isinstance(count_val, (int, float)) else str(count_val)
        missing_str = f"{missing_count_val:,}" if isinstance(missing_count_val, (int, float)) else str(missing_count_val)
        unique_str = f"{unique_count_val:,}" if isinstance(unique_count_val, (int, float)) else str(unique_count_val)
        
        basic_stats = f"""🏷️ 基础统计 (Basic Statistics)

计数 (Count): {count_str}
缺失 (Missing): {missing_str} ({analysis.get('missing_percentage', 0):.1f}%)
唯一值 (Unique): {unique_str}
基数 (Cardinality): {analysis.get('cardinality', 0):.4f}
众数 (Mode): {analysis.get('mode', 'N/A')}
高基数 (High Card): {'是 (Yes)' if analysis.get('is_high_cardinality', False) else '否 (No)'}"""
        
        basic_frame.text = basic_stats
        basic_frame.paragraphs[0].font.size = Pt(14)
        basic_frame.paragraphs[0].font.bold = True
        for i in range(1, len(basic_frame.paragraphs)):
            basic_frame.paragraphs[i].font.size = Pt(11)
        
        # Top Values Box (Right)
        value_counts = analysis.get('value_counts', {})
        value_percentages = analysis.get('value_percentages', {})
        
        if value_counts:
            values_box = slide.shapes.add_textbox(Inches(6.7), Inches(1.3), Inches(5.8), Inches(2.5))
            values_frame = values_box.text_frame
            values_frame.word_wrap = True
            
            values_text = "🔝 高频值 (Top Values)\n\n"
            for value, count in list(value_counts.items())[:6]:  # Top 6
                percentage = value_percentages.get(value, 0)
                value_str = str(value)[:20] + "..." if len(str(value)) > 20 else str(value)
                count_str = f"{count:,}" if isinstance(count, (int, float)) else str(count)
                values_text += f"{value_str}: {count_str} ({percentage:.1f}%)\n"
            
            values_frame.text = values_text
            values_frame.paragraphs[0].font.size = Pt(14)
            values_frame.paragraphs[0].font.bold = True
            for i in range(1, len(values_frame.paragraphs)):
                values_frame.paragraphs[i].font.size = Pt(10)
        
        # Rare Categories Box (Bottom Left)
        rare_categories = analysis.get('rare_categories', [])
        if rare_categories:
            rare_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(5.8), Inches(1.5))
            rare_frame = rare_box.text_frame
            rare_frame.word_wrap = True
            
            rare_text = f"🔍 稀有类别 (Rare Categories) - {len(rare_categories)} 个:\n"
            rare_display = rare_categories[:8] if len(rare_categories) > 8 else rare_categories
            rare_text += ", ".join(str(cat)[:15] for cat in rare_display)
            if len(rare_categories) > 8:
                rare_text += f"... (+{len(rare_categories) - 8} 更多)"
            
            rare_frame.text = rare_text
            rare_frame.paragraphs[0].font.size = Pt(12)
            rare_frame.paragraphs[0].font.bold = True
        
        # Recommendations Box (Bottom)
        recommendations = analysis.get('recommendations', [])
        if recommendations:
            rec_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(12), Inches(1.3))
            rec_frame = rec_box.text_frame
            rec_frame.word_wrap = True
            
            rec_text = "💡 清洗建议 (Recommendations):\n"
            for i, rec in enumerate(recommendations[:4], 1):
                rec_text += f"{i}. {rec}\n"
            
            rec_frame.text = rec_text
            rec_frame.paragraphs[0].font.size = Pt(14)
            rec_frame.paragraphs[0].font.bold = True
            for i in range(1, len(rec_frame.paragraphs)):
                rec_frame.paragraphs[i].font.size = Pt(10)
    
    def _create_datetime_stats_layout(self, slide, analysis: Dict[str, Any]):
        """Create beautiful layout for datetime statistics."""
        # Basic Stats Box (Left)
        basic_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(5.8), Inches(2.5))
        basic_frame = basic_box.text_frame
        basic_frame.word_wrap = True
        
        count_val = analysis.get('count', 0)
        missing_count_val = analysis.get('missing_count', 0)
        
        count_str = f"{count_val:,}" if isinstance(count_val, (int, float)) else str(count_val)
        missing_str = f"{missing_count_val:,}" if isinstance(missing_count_val, (int, float)) else str(missing_count_val)
        
        basic_stats = f"""📅 基础统计 (Basic Statistics)

计数 (Count): {count_str}
缺失 (Missing): {missing_str} ({analysis.get('missing_percentage', 0):.1f}%)
最早日期 (Min Date): {analysis.get('min_date', 'N/A')}
最晚日期 (Max Date): {analysis.get('max_date', 'N/A')}
日期范围 (Range): {analysis.get('date_range', 'N/A')}
年份跨度 (Year Range): {analysis.get('year_range', 'N/A')} 年"""
        
        basic_frame.text = basic_stats
        basic_frame.paragraphs[0].font.size = Pt(14)
        basic_frame.paragraphs[0].font.bold = True
        for i in range(1, len(basic_frame.paragraphs)):
            basic_frame.paragraphs[i].font.size = Pt(11)
        
        # Time Patterns Box (Right)
        patterns_box = slide.shapes.add_textbox(Inches(6.7), Inches(1.3), Inches(5.8), Inches(2.5))
        patterns_frame = patterns_box.text_frame
        patterns_frame.word_wrap = True
        
        unique_years_val = analysis.get('unique_years', 0)
        unique_months_val = analysis.get('unique_months', 0)
        unique_days_val = analysis.get('unique_days', 0)
        
        years_str = f"{unique_years_val:,}" if isinstance(unique_years_val, (int, float)) else str(unique_years_val)
        months_str = f"{unique_months_val:,}" if isinstance(unique_months_val, (int, float)) else str(unique_months_val)
        days_str = f"{unique_days_val:,}" if isinstance(unique_days_val, (int, float)) else str(unique_days_val)
        
        patterns_text = f"""🕒 时间模式 (Time Patterns)

唯一年份 (Unique Years): {years_str}
唯一月份 (Unique Months): {months_str}
唯一天数 (Unique Days): {days_str}"""
        
        patterns_frame.text = patterns_text
        patterns_frame.paragraphs[0].font.size = Pt(14)
        patterns_frame.paragraphs[0].font.bold = True
        for i in range(1, len(patterns_frame.paragraphs)):
            patterns_frame.paragraphs[i].font.size = Pt(11)
        
        # Distribution Boxes (Bottom)
        weekday_dist = analysis.get('weekday_distribution', {})
        if weekday_dist:
            weekday_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(5.8), Inches(1.5))
            weekday_frame = weekday_box.text_frame
            weekday_frame.word_wrap = True
            
            weekday_text = "📅 星期分布 (Weekday Distribution):\n"
            for day, count in list(weekday_dist.items())[:7]:
                count_str = f"{count:,}" if isinstance(count, (int, float)) else str(count)
                weekday_text += f"{day}: {count_str}  "
            
            weekday_frame.text = weekday_text
            weekday_frame.paragraphs[0].font.size = Pt(12)
            weekday_frame.paragraphs[0].font.bold = True
        
        # Recommendations Box
        recommendations = analysis.get('recommendations', [])
        if recommendations:
            rec_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(12), Inches(1.3))
            rec_frame = rec_box.text_frame
            rec_frame.word_wrap = True
            
            rec_text = "💡 清洗建议 (Recommendations):\n"
            for i, rec in enumerate(recommendations[:4], 1):
                rec_text += f"{i}. {rec}\n"
            
            rec_frame.text = rec_text
            rec_frame.paragraphs[0].font.size = Pt(14)
            rec_frame.paragraphs[0].font.bold = True
            for i in range(1, len(rec_frame.paragraphs)):
                rec_frame.paragraphs[i].font.size = Pt(10)
    
    def _create_unknown_stats_layout(self, slide, analysis: Dict[str, Any]):
        """Create layout for unknown data types."""
        # Basic Info Box
        basic_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12), Inches(2))
        basic_frame = basic_box.text_frame
        basic_frame.word_wrap = True
        
        count_val = analysis.get('count', 0)
        missing_count_val = analysis.get('missing_count', 0)
        
        count_str = f"{count_val:,}" if isinstance(count_val, (int, float)) else str(count_val)
        missing_str = f"{missing_count_val:,}" if isinstance(missing_count_val, (int, float)) else str(missing_count_val)
        
        basic_text = f"""❓ 数据类型未知 (Unknown Data Type)

类型 (Type): {analysis.get('type', 'Unknown')}
计数 (Count): {count_str}
缺失 (Missing): {missing_str} ({analysis.get('missing_percentage', 0):.1f}%)

该列的数据类型无法自动识别，建议手动检查数据内容。"""
        
        basic_frame.text = basic_text
        basic_frame.paragraphs[0].font.size = Pt(16)
        basic_frame.paragraphs[0].font.bold = True
        for i in range(1, len(basic_frame.paragraphs)):
            basic_frame.paragraphs[i].font.size = Pt(12)
        
        # Recommendations
        recommendations = analysis.get('recommendations', [])
        if recommendations:
            rec_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(12), Inches(2))
            rec_frame = rec_box.text_frame
            rec_frame.word_wrap = True
            
            rec_text = "💡 清洗建议 (Recommendations):\n"
            for i, rec in enumerate(recommendations[:5], 1):
                rec_text += f"{i}. {rec}\n"
            
            rec_frame.text = rec_text
            rec_frame.paragraphs[0].font.size = Pt(14)
            rec_frame.paragraphs[0].font.bold = True
            for i in range(1, len(rec_frame.paragraphs)):
                rec_frame.paragraphs[i].font.size = Pt(11)
    
    def _create_data_conversion_slide(self, prs: Presentation, results: Dict[str, Any]):
        """Create data type conversion analysis slide with beautiful table UI."""
        slide_layout = prs.slide_layouts[6]  # Blank layout for custom design
        slide = prs.slides.add_slide(slide_layout)
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "🔄 数据类型转换分析 (Data Type Conversion Analysis)"
        title_frame.paragraphs[0].font.size = Pt(24)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
        
        # Get conversion data
        conversion_data = results.get('data_conversion', {})
        conversion_report = conversion_data.get('conversion_report', {})
        conversion_analysis = conversion_data.get('conversion_analysis', {})
        
        # Statistics summary box (top section)
        convertible_count = len(conversion_report.get('convertible_columns', []))
        partially_convertible_count = len(conversion_report.get('partially_convertible_columns', {}))
        unconvertible_count = len(conversion_report.get('unconvertible_columns', []))
        already_float_count = len(conversion_report.get('already_float_columns', []))
        total_columns = conversion_analysis.get('total_columns', 0)
        success_rate = conversion_analysis.get('conversion_success_rate', 0)
        
        # Create beautiful stats boxes (top row)
        stats_boxes = [
            ("完全可转换", convertible_count, "🟢", RGBColor(34, 139, 34)),
            ("部分可转换", partially_convertible_count, "🟡", RGBColor(255, 165, 0)),
            ("无法转换", unconvertible_count, "🔴", RGBColor(220, 20, 60)),
            ("已为数值型", already_float_count, "✅", RGBColor(70, 130, 180))
        ]
        
        box_width = 2.8
        box_height = 1.2
        start_x = 0.7
        
        for i, (label, count, emoji, color) in enumerate(stats_boxes):
            x_pos = start_x + i * (box_width + 0.2)
            
            # Create colored background box
            stats_box = slide.shapes.add_textbox(Inches(x_pos), Inches(1.3), Inches(box_width), Inches(box_height))
            stats_frame = stats_box.text_frame
            stats_frame.text = f"{emoji} {label}\n{count} 列"
            
            # Style the box
            stats_frame.paragraphs[0].font.size = Pt(14)
            stats_frame.paragraphs[0].font.bold = True
            stats_frame.paragraphs[0].font.color.rgb = color
            stats_frame.paragraphs[0].alignment = 1  # Center alignment
            
            # Add border effect
            stats_box.fill.solid()
            stats_box.fill.fore_color.rgb = RGBColor(248, 248, 248)
            stats_box.line.color.rgb = color
            stats_box.line.width = Pt(2)
        
        # Success rate box (center)
        success_box = slide.shapes.add_textbox(Inches(4), Inches(2.8), Inches(4.5), Inches(1))
        success_frame = success_box.text_frame
        success_frame.text = f"🎯 整体转换成功率: {success_rate:.1%}\n({convertible_count + partially_convertible_count}/{total_columns} 列可转换)"
        success_frame.paragraphs[0].font.size = Pt(16)
        success_frame.paragraphs[0].font.bold = True
        success_frame.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)
        success_frame.paragraphs[0].alignment = 1  # Center
        
        # Convertible columns list (left side)
        if conversion_report.get('convertible_columns'):
            conv_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(5.8), Inches(2.8))
            conv_frame = conv_box.text_frame
            conv_frame.word_wrap = True
            
            conv_text = "🟢 完全可转换列 (Fully Convertible):\n\n"
            cols = conversion_report['convertible_columns'][:8]  # Show first 8
            for col in cols:
                conv_text += f"• {col}\n"
            if len(conversion_report['convertible_columns']) > 8:
                conv_text += f"... 及其他 {len(conversion_report['convertible_columns']) - 8} 列"
            
            conv_frame.text = conv_text
            conv_frame.paragraphs[0].font.size = Pt(14)
            conv_frame.paragraphs[0].font.bold = True
            conv_frame.paragraphs[0].font.color.rgb = RGBColor(34, 139, 34)
            
            # Style list items
            for i in range(1, len(conv_frame.paragraphs)):
                conv_frame.paragraphs[i].font.size = Pt(11)
                conv_frame.paragraphs[i].font.color.rgb = RGBColor(60, 60, 60)
        
        # Recommendations box (right side)
        recommendations = conversion_analysis.get('recommendations', [])
        if recommendations:
            rec_box = slide.shapes.add_textbox(Inches(6.7), Inches(4.2), Inches(5.8), Inches(2.8))
            rec_frame = rec_box.text_frame
            rec_frame.word_wrap = True
            
            rec_text = "💡 建议 (Recommendations):\n\n"
            for i, rec in enumerate(recommendations[:4], 1):
                rec_text += f"{i}. {rec}\n"
            
            rec_frame.text = rec_text
            rec_frame.paragraphs[0].font.size = Pt(14)
            rec_frame.paragraphs[0].font.bold = True
            rec_frame.paragraphs[0].font.color.rgb = RGBColor(255, 140, 0)
            
            # Style list items
            for i in range(1, len(rec_frame.paragraphs)):
                rec_frame.paragraphs[i].font.size = Pt(11)
                rec_frame.paragraphs[i].font.color.rgb = RGBColor(60, 60, 60)